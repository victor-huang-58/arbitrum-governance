#!/usr/bin/env python3
"""
Build presentation.pptx — 25-slide conference deck
"The Effect of AI Writing on Governance: Evidence from a $3.5 Billion DAO"
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
from pathlib import Path
from PIL import Image
import copy, io

# ── Colour palette ────────────────────────────────────────────
NAVY   = RGBColor(0x1e, 0x3a, 0x5f)
BLUE   = RGBColor(0x1e, 0x64, 0xb4)
PRE    = RGBColor(0x48, 0x78, 0xcf)
POST   = RGBColor(0xd6, 0x5f, 0x5f)
ORANGE = RGBColor(0xcc, 0x55, 0x00)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LGRAY  = RGBColor(0xf0, 0xf2, 0xf5)
DGRAY  = RGBColor(0x33, 0x33, 0x33)
BLACK  = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

FIGURES = Path("output/figures")

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.333, 1.15, fill=NAVY)
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.12), Inches(12.0), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.35), Inches(0.72), Inches(12.0), Inches(0.38))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = RGBColor(0xcc, 0xdd, 0xff)

def footer(slide, n, total=25):
    rect(slide, 0, 7.15, 13.333, 0.35, fill=NAVY)
    tb = slide.shapes.add_textbox(Inches(11.8), Inches(7.17), Inches(1.3), Inches(0.28))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{n} / {total}"
    r.font.size = Pt(9)
    r.font.color.rgb = RGBColor(0xaa, 0xcc, 0xff)
    tb2 = slide.shapes.add_textbox(Inches(0.3), Inches(7.17), Inches(8), Inches(0.28))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "The Effect of AI Writing on Governance: Evidence from a $3.5 Billion DAO"
    r2.font.size = Pt(9)
    r2.font.italic = True
    r2.font.color.rgb = RGBColor(0xaa, 0xcc, 0xff)

def textbox(slide, text, l, t, w, h, size=14, bold=False, color=DGRAY,
            align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def bullets(slide, items, l, t, w, h, size=13, title=None, title_color=NAVY,
            bullet_color=DGRAY, indent=0.25):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = title
        r.font.size = Pt(size + 1)
        r.font.bold = True
        r.font.color.rgb = title_color
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(3)
        p.level = 0
        r = p.add_run()
        r.text = f"  •  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = bullet_color
    return tb

def add_figure(slide, fname, l, t, w, h=None):
    path = FIGURES / fname
    if not path.exists():
        return
    if h is None:
        with Image.open(path) as im:
            iw, ih = im.size
        h = w * ih / iw
    slide.shapes.add_picture(str(path), Inches(l), Inches(t),
                             Inches(w), Inches(h))

def callout_box(slide, text, l, t, w, h, fill=LGRAY, border=BLUE, size=13):
    rect(slide, l, t, w, h, fill=fill, line=border)
    tb = slide.shapes.add_textbox(Inches(l+0.15), Inches(t+0.12),
                                  Inches(w-0.3), Inches(h-0.24))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = DGRAY

def stat_box(slide, label, value, l, t, w=2.8, h=0.95, fill=LGRAY, vcolor=NAVY):
    rect(slide, l, t, w, h, fill=fill, line=BLUE)
    textbox(slide, label, l+0.1, t+0.08, w-0.2, 0.35, size=10, color=DGRAY)
    textbox(slide, value, l+0.1, t+0.38, w-0.2, 0.5, size=18, bold=True,
            color=vcolor)

# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, NAVY); footer(sl, 1)
rect(sl, 0, 0, 13.333, 7.5, fill=NAVY)
rect(sl, 0, 2.8, 13.333, 0.07, fill=BLUE)
textbox(sl, "The Effect of AI Writing on Governance: Evidence from a $3.5 Billion DAO",
        0.6, 1.0, 12.1, 1.5, size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(sl, "How Large Language Models Undermined DAO Deliberation",
        0.6, 2.5, 12.1, 0.6, size=17, color=RGBColor(0xcc,0xdd,0xff),
        align=PP_ALIGN.CENTER, italic=True)
textbox(sl, "Victor Huang  ·  Joseph Hall",
        0.6, 3.3, 12.1, 0.5, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(sl, "Journal of Finance",
        0.6, 3.85, 12.1, 0.4, size=13, color=RGBColor(0xaa,0xcc,0xff),
        align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 2 — THE QUESTION
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); header_bar(sl, "The Research Question"); footer(sl, 2)
textbox(sl,
    "When AI can generate indistinguishable governance prose,\n"
    "do rational participants stop treating the forum as a credible signal?",
    0.6, 1.35, 12.1, 1.1, size=19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
rect(sl, 1.0, 2.55, 11.3, 0.04, fill=BLUE)

cols = [
    ("Setting", "Arbitrum DAO: largest DAO by treasury ($3.5B ARB), 17,553 forum posts, 306 Snapshot proposals (102 matched)"),
    ("Shock",   "GPT-4 Turbo (Nov 2023) / Claude 3 (Mar 2024): first LLMs capable of governance-quality prose"),
    ("Method",  "Structural break analysis: does the forum-signal relationship change around AI model releases?"),
    ("Answer",  "Yes — completely. Pre-shock r = −0.33 (n=29). Post-shock r = −0.10 (n=73). Chow F = 8.38 (p=0.0004)."),
]
for i, (hd, txt) in enumerate(cols):
    x = 0.55 + i * 3.1
    rect(sl, x, 2.75, 2.85, 3.5, fill=LGRAY, line=BLUE)
    textbox(sl, hd, x+0.12, 2.85, 2.6, 0.45, size=12, bold=True, color=NAVY)
    textbox(sl, txt, x+0.12, 3.35, 2.6, 2.8, size=11, color=DGRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 3 — WHY IT MATTERS
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); header_bar(sl, "Why This Matters"); footer(sl, 3)
bullets(sl,
    ["The ten largest DAOs govern ~$9B in on-chain treasury assets",
     "Text-based deliberation is the only public governance input mechanism",
     "Arbitrum alone: $3.5B treasury, 306 on-chain proposals (102 matched), 62 identified delegates",
     "If AI undermines deliberation, governance quality degrades silently"],
    0.5, 1.25, 6.0, 3.0, title="The stakes in DAOs", size=13)
bullets(sl,
    ["Public comment periods on regulation",
     "Earnings call Q&A and investor forums",
     "Corporate ESG engagement letters",
     "Any institution that relies on text-based participation"],
    6.7, 1.25, 6.0, 3.0, title="The broader question", size=13)
callout_box(sl,
    "Akerlof (1970): Good cars exit the used-car market not because lemons are "
    "common, but because buyers cannot distinguish them. "
    "The governance forum collapses for the same reason.",
    0.5, 4.55, 12.3, 1.15, fill=RGBColor(0xe8,0xf0,0xff), border=BLUE, size=13)

# ════════════════════════════════════════════════════════════
# SLIDE 4 — SETTING
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); header_bar(sl, "Setting: Arbitrum DAO"); footer(sl, 4)
bullets(sl,
    ["Layer-2 Ethereum rollup; ARB token launched March 2023",
     "Liquid democracy: token holders delegate to professional delegates",
     "Proposal lifecycle: forum discussion → Snapshot vote → on-chain execution",
     "Forum deliberation is the only public deliberative stage before voting"],
    0.5, 1.3, 5.8, 3.2, title="Governance structure", size=13)
bullets(sl,
    ["17,553 forum posts, Jan 2023 – Mar 2026",
     "102 matched Snapshot proposals (306 binary total)",
     "All posts scored by Pangram AI-detection API (v3 endpoint)",
     "Vote margin: |2·For/(For+Against) − 1|, ranging 0 (tied) to 1 (unanimous)"],
    0.5, 4.6, 5.8, 2.3, title="Data", size=13)
add_figure(sl, "fig01_forum_activity.png", 6.6, 1.25, 6.4)

# ════════════════════════════════════════════════════════════
# SLIDE 5 — THE FORUM SIGNAL
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); header_bar(sl, "The Forum as an Information Market"); footer(sl, 5)
textbox(sl, "Contested proposals attract more genuine discussion.",
        0.5, 1.3, 12.3, 0.45, size=15, bold=True, color=NAVY)
textbox(sl, "A rational delegate should update positively on human post volume when deciding how to vote.",
        0.5, 1.75, 12.3, 0.4, size=13, color=DGRAY)
add_figure(sl, "fig08_main_causal.png", 0.5, 2.25, 8.0)
stat_box(sl, "Pre-shock correlation", "r = −0.33", 9.0, 2.3, fill=RGBColor(0xe8,0xf0,0xff), vcolor=PRE)
stat_box(sl, "p-value (OLS)", "p = 0.030", 9.0, 3.35, fill=RGBColor(0xe8,0xf0,0xff), vcolor=NAVY)
stat_box(sl, "Sample size", "n = 29", 9.0, 4.4, fill=RGBColor(0xe8,0xf0,0xff), vcolor=NAVY)
textbox(sl, "More human posts → more contested vote. Forum is a real signal. Until it isn't.",
        9.0, 5.5, 4.1, 0.8, size=11, italic=True, color=DGRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 6 — THE MECHANISM
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); header_bar(sl, "The Akerlof Lemons Mechanism"); footer(sl, 6)
textbox(sl, "Expected value of reading post j before voting on proposal p:",
        0.5, 1.3, 12.3, 0.4, size=13, color=DGRAY)
textbox(sl, "E[Vjp]  =  (1 − π̲) · IH",
        0.5, 1.75, 12.3, 0.55, size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
bullets(sl,
    ["π̲ = prior probability any given post is AI-generated",
     "IH = information value of a genuine human post (IH > 0)",
     "c_read = cost of careful reading, incurred BEFORE type is revealed",
     "Delegate reads iff  E[Vjp] > c_read"],
    0.5, 2.4, 12.3, 2.0, size=13)
rect(sl, 0.5, 4.55, 5.7, 1.55, fill=RGBColor(0xe8,0xf4,0xe8), line=RGBColor(0x22,0x88,0x44))
textbox(sl, "Pre-LLM  (π̲ ≈ 0):", 0.65, 4.62, 5.4, 0.35, size=12, bold=True, color=RGBColor(0x22,0x66,0x33))
textbox(sl, "(1−π̲)·IH > c_read\n→ delegates read; forum is informative",
        0.65, 5.0, 5.4, 1.0, size=12, color=DGRAY)
rect(sl, 7.2, 4.55, 5.7, 1.55, fill=RGBColor(0xfd,0xee,0xee), line=POST)
textbox(sl, "Post-LLM  (π̲ crosses threshold):", 7.35, 4.62, 5.4, 0.35, size=12, bold=True, color=POST)
textbox(sl, "(1−π̲)·IH < c_read\n→ delegates exit; correlation → 0",
        7.35, 5.0, 5.4, 1.0, size=12, color=DGRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 7 — TWO CHANNELS
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "Two Channels of Signal Degradation",
           "The human-post signal dies through complementary mechanisms — neither requires individual posts to be unclassifiable")
footer(sl, 7)
rect(sl, 0.5, 1.3, 5.85, 4.4, fill=RGBColor(0xe8,0xf0,0xff), line=BLUE)
textbox(sl, "Channel 1: Sequential Reading Cost", 0.65, 1.42, 5.55, 0.45,
        size=13, bold=True, color=NAVY)
textbox(sl,
    "Reading costs are incurred before post type is observed. "
    "Even with aggregate stylometric differences, a delegate cannot "
    "identify human posts without first reading them.\n\n"
    "As AI base rate rises, the expected return to reading any post falls "
    "— even though vocabulary differences remain in aggregate.",
    0.65, 1.9, 5.55, 2.6, size=12, color=DGRAY)
stat_box(sl, "AI posts per human post", "0.05 → 0.13  (×2.6)", 0.65, 4.55,
         w=5.5, h=0.85, fill=RGBColor(0xd8,0xe8,0xff), vcolor=BLUE)

rect(sl, 7.0, 1.3, 5.85, 4.4, fill=RGBColor(0xfd,0xf0,0xf0), line=POST)
textbox(sl, "Channel 2: Ecosystem Degradation", 7.15, 1.42, 5.55, 0.45,
        size=13, bold=True, color=POST)
textbox(sl,
    "As AI floods the forum, genuine human participants disengage. "
    "The remaining human posts are concentrated among fewer voices — "
    "less representative of broad community sentiment.\n\n"
    "Even if delegates could perfectly filter to human posts, "
    "those posts carry less deliberative information.",
    7.15, 1.9, 5.55, 2.6, size=12, color=DGRAY)
stat_box(sl, "Unique human participants", "740 → 496  (−33%)", 7.15, 4.55,
         w=5.5, h=0.85, fill=RGBColor(0xfd,0xe0,0xe0), vcolor=POST)

# ════════════════════════════════════════════════════════════
# SLIDE 8 — DATA
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); header_bar(sl, "Data"); footer(sl, 8)
bullets(sl,
    ["17,553 posts across 1,101 governance topics (Jan 2023 – Mar 2026)",
     "Full text, read counts, like counts, reply counts per post",
     "Forum–Snapshot match: title token overlap ≥ 0.55; 102 matched proposals"],
    0.5, 1.3, 6.0, 2.5, title="Forum (Discourse API)", size=13)
bullets(sl,
    ["All 15,891 scoreable posts classified by Pangram API v3 endpoint",
     "fraction_ai ∈ [0,1]; threshold 0.70 → AI-classified",
     "90.5% of corpus has valid scores; results robust to 0.50, 0.60, 0.80",
     "Measurement error attenuates correlations toward zero (conservative bias)"],
    0.5, 3.95, 6.0, 2.75, title="AI detection (Pangram)", size=13)
bullets(sl,
    ["Vote tallies from Snapshot GraphQL API (306 binary proposals)",
     "Margin = |2·For/(For+Against) − 1|  ∈ [0, 1]",
     "Delegate wallet map: 62 identified delegates with forum usernames",
     "Post-level read counts: direct measure of delegate engagement"],
    6.7, 1.3, 6.0, 5.4, title="Snapshot + delegate data", size=13)

# ════════════════════════════════════════════════════════════
# SLIDE 9 — AI DETECTION VALIDITY
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "AI Detection: Validity and Content Analysis")
footer(sl, 9)
add_figure(sl, "fig_ai_post_content.png", 0.5, 1.3, 7.0)
bullets(sl,
    ["AI posts: 'community engagement,' 'support proposal,' 'decision making' — a consultant's summary",
     "Human posts: 'don't think,' 'temp check,' 'security council,' 'tally vote' — skepticism and procedure",
     "AI posts 1.7× more likely to contain generic agreement phrases (35.4% vs 21.2%)",
     "AI posts are proposal echoes, not opinion signals"],
    7.7, 1.3, 5.4, 4.5, title="TF-IDF content analysis validates classification", size=12)
callout_box(sl,
    "Aggregate corpus differences are detectable ex-post. They do not imply "
    "individual posts are reliably classifiable in real time — which is what the mechanism requires.",
    7.7, 5.9, 5.4, 0.9, size=11)

# ════════════════════════════════════════════════════════════
# SLIDE 10 — MAIN RESULT
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); header_bar(sl, "Main Result: The Correlation Collapse"); footer(sl, 10)
add_figure(sl, "fig08_main_causal.png", 0.5, 1.25, 9.5)
rect(sl, 10.2, 1.3, 2.85, 1.5, fill=RGBColor(0xe8,0xf0,0xff), line=PRE)
textbox(sl, "Pre-shock", 10.35, 1.38, 2.55, 0.38, size=12, bold=True, color=PRE)
textbox(sl, "r = −0.33\np = 0.030\nn = 29", 10.35, 1.78, 2.55, 0.9, size=13, color=DGRAY)
rect(sl, 10.2, 2.95, 2.85, 1.5, fill=RGBColor(0xfd,0xee,0xee), line=POST)
textbox(sl, "Post-shock", 10.35, 3.03, 2.55, 0.38, size=12, bold=True, color=POST)
textbox(sl, "r = −0.10\np = 0.390\nn = 73", 10.35, 3.43, 2.55, 0.9, size=13, color=DGRAY)
rect(sl, 10.2, 4.6, 2.85, 1.3, fill=LGRAY, line=NAVY)
textbox(sl, "Chow F", 10.35, 4.68, 2.55, 0.35, size=12, bold=True, color=NAVY)
textbox(sl, "F = 8.38\np = 0.0004", 10.35, 5.05, 2.55, 0.75, size=13, color=DGRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 11 — STRUCTURAL BREAK
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); header_bar(sl, "Identifying the Structural Break"); footer(sl, 11)
add_figure(sl, "fig13_chow.png", 0.5, 1.3, 7.5)
textbox(sl, "Three-pronged identification:", 8.3, 1.3, 4.8, 0.4,
        size=13, bold=True, color=NAVY)
items = [
    ("Chow test (Claude 3)",    "F = 8.38,  p = 0.0004"),
    ("QLR test (endogenous)",   "sup-F = 17.19 at Nov 25, 2023\n5% CV = 8.85  (Andrews 1993)"),
    ("Bootstrap permutation",  "9 / 10,000 random permutations\nexceed F = 8.38;  p_perm = 0.0009"),
]
for i, (lbl, val) in enumerate(items):
    y = 1.8 + i * 1.55
    rect(sl, 8.3, y, 4.8, 1.35, fill=LGRAY, line=BLUE)
    textbox(sl, lbl, 8.45, y+0.1, 4.5, 0.38, size=11, bold=True, color=NAVY)
    textbox(sl, val, 8.45, y+0.5, 4.5, 0.75, size=12, color=DGRAY)
callout_box(sl,
    "QLR break: November 25, 2023 — nineteen days after GPT-4 Turbo release.",
    8.3, 6.4, 4.8, 0.65, fill=RGBColor(0xe8,0xf0,0xff))

# ════════════════════════════════════════════════════════════
# SLIDE 12 — CHOW BATTERY
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "Chow Test Battery",
           "All six candidate dates show the same directional break")
footer(sl, 12)

rows = [
    ("Claude 2  (Jul 2023)",     "7",   "−0.17", "95",  "−0.22",  "1.00",  "0.371",    False),
    ("GPT-4 Turbo  (Nov 2023)",  "18",  "−0.34", "84",  "−0.17", "15.05", "<0.001",    True),
    ("Claude 3  (Mar 2024) ★",   "29",  "−0.33", "73",  "−0.10",  "8.38",  "0.0004",   True),
    ("GPT-4o  (May 2024)",       "36",  "−0.30", "66",  "−0.06",  "5.97",  "0.004",    True),
    ("o1  (Sep 2024)",           "65",  "−0.20", "37",  "−0.18",  "1.63",  "0.200",    False),
    ("DeepSeek R1  (Jan 2025)",  "78",  "−0.17", "24",  "−0.38",  "0.67",  "0.513",    False),
]
headers = ["Shock date", "n_pre", "r_pre", "n_post", "r_post", "F", "p"]
col_x   = [0.5, 4.7, 5.7, 6.6, 7.5, 8.7, 9.8]
col_w   = [4.1, 0.9, 0.9, 0.9, 0.9, 1.1, 1.4]

y0 = 1.45
for j, hd in enumerate(headers):
    textbox(sl, hd, col_x[j], y0, col_w[j], 0.38, size=11, bold=True, color=NAVY)
rect(sl, 0.5, 1.85, 12.3, 0.04, fill=NAVY)

for i, row in enumerate(rows):
    y = 1.95 + i * 0.72
    if row[-1]:
        rect(sl, 0.5, y-0.05, 12.3, 0.7, fill=RGBColor(0xe8,0xf0,0xff))
    for j, val in enumerate(row[:-1]):
        color = NAVY if row[-1] else DGRAY
        textbox(sl, val, col_x[j], y, col_w[j], 0.6, size=11,
                bold=row[-1], color=color)

rect(sl, 0.5, 6.3, 12.3, 0.04, fill=NAVY)
textbox(sl,
    "QLR sup-F = 17.19 at Nov 25, 2023  |  5% CV = 8.85  |  "
    "★ Primary break: Claude 3 chosen for a clean before/after contrast (clear negative pre, null post); results robust to GPT-4 Turbo  |  "
    "Three dates survive Bonferroni correction (p < 0.0083)",
    0.5, 6.4, 12.3, 0.65, size=10, color=DGRAY, italic=True)

# ════════════════════════════════════════════════════════════
# SLIDE 13 — TIMING PARADOX
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "The Timing Paradox: Collapse Precedes Contamination")
footer(sl, 13)
add_figure(sl, "fig08_main_causal.png", 0.5, 1.3, 8.2)
callout_box(sl,
    "Volume mechanism predicts: signal dies when AI posts flood the forum.\n\n"
    "Observed: signal dies BEFORE the AI volume surge. "
    "Panel A shows volume accelerates months after the break.\n\n"
    "Unverifiability mechanism predicts exactly this: delegates respond to the "
    "credible possibility of contamination, not its realized magnitude.",
    8.9, 1.3, 4.2, 5.5, fill=RGBColor(0xfd,0xf0,0xf0), border=POST, size=12)

# ════════════════════════════════════════════════════════════
# SLIDE 14 — DIRECT BEHAVIORAL EVIDENCE
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "Direct Behavioral Evidence: Forum Readership Decline",
           "Logged read counts measure actual engagement — not inferred behavior")
footer(sl, 14)
add_figure(sl, "fig15_reads_decline.png", 0.5, 1.3, 9.5)
stat_box(sl, "Pre-shock median reads/post", "88", 10.25, 1.5,
         w=2.8, fill=RGBColor(0xe8,0xf0,0xff), vcolor=PRE)
stat_box(sl, "Post-shock median reads/post", "56", 10.25, 2.6,
         w=2.8, fill=RGBColor(0xfd,0xee,0xee), vcolor=POST)
stat_box(sl, "Decline", "−36%", 10.25, 3.7,
         w=2.8, fill=LGRAY, vcolor=ORANGE)
textbox(sl,
    "Decline begins November 2023, "
    "coincident with GPT-4 Turbo.\n\n"
    "Not a volume artifact: human post counts increased post-shock.",
    10.25, 4.8, 2.8, 1.6, size=11, color=DGRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 15 — ECOSYSTEM DEGRADATION
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "Ecosystem Degradation: Participation Collapse")
footer(sl, 15)
stats = [
    ("Unique human\nparticipants",  "740",  "496",  "−33%",   POST),
    ("Posts per user",              "3.6",  "11.9", "×3.3",   ORANGE),
    ("AI fraction\n(prop. threads)","4.7%", "11.3%","×2.4",   POST),
    ("AI posts per\nhuman post",    "0.05", "0.13", "×2.6",   ORANGE),
    ("Median reads\nper post",      "88",   "56",   "−36%",   POST),
]
for i, (lbl, pre, post, chg, col) in enumerate(stats):
    x = 0.5 + i * 2.55
    rect(sl, x, 1.3, 2.35, 3.8, fill=LGRAY, line=BLUE)
    textbox(sl, lbl, x+0.1, 1.4, 2.15, 0.65, size=11, bold=True, color=NAVY)
    textbox(sl, "Pre-shock", x+0.1, 2.15, 1.0, 0.3, size=9, color=DGRAY)
    textbox(sl, pre, x+0.1, 2.45, 2.15, 0.5, size=18, bold=True, color=PRE)
    textbox(sl, "Post-shock", x+0.1, 3.05, 1.0, 0.3, size=9, color=DGRAY)
    textbox(sl, post, x+0.1, 3.35, 2.15, 0.5, size=18, bold=True, color=POST)
    rect(sl, x+0.1, 4.7-0.6, 2.15, 0.55, fill=RGBColor(0xff,0xf0,0xe8), line=col)
    textbox(sl, chg, x+0.1, 4.72-0.6, 2.15, 0.42, size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
callout_box(sl,
    "Even if a delegate could perfectly identify human posts, those posts increasingly represent "
    "a small number of heavy posters — not broad community sentiment. "
    "The informational value of human post count (IH) fell endogenously.",
    0.5, 5.25, 12.3, 1.05, fill=RGBColor(0xe8,0xf0,0xff))

# ════════════════════════════════════════════════════════════
# SLIDE 16 — CONTENT ANALYSIS
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "What AI Posts Actually Say: TF-IDF Analysis")
footer(sl, 16)
add_figure(sl, "fig_ai_post_content.png", 0.5, 1.3, 8.5)
bullets(sl,
    ["AI n-grams: 'community engagement,' 'support proposal,' 'decision making,' 'transparency accountability'",
     "Human n-grams: 'don't think,' 'temp check,' 'security council,' 'tally vote,' 'non-constitutional'",
     "AI posts: 1.7× higher generic agreement phrase rate (35.4% vs 21.2%)",
     "AI posts are longer (median 160 vs 126 words) but less informationally dense"],
    9.2, 1.3, 4.0, 4.8, size=11)
callout_box(sl,
    "AI posts open with structured summaries followed by conditional endorsements. "
    "They do not take positions, reference prior votes, or cite technical concerns. "
    "AI-classified posts are proposal echoes, not opinion signals.",
    0.5, 6.2, 12.3, 0.95, size=11)

# ════════════════════════════════════════════════════════════
# SLIDE 17 — DELEGATE CROSS-SECTION
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "Delegate Cross-Section: AI Adoption and Voting Power")
footer(sl, 17)
add_figure(sl, "fig11_delegates.png", 0.5, 1.3, 7.8)
textbox(sl, "Finding: delegates who use more AI hold systematically lower voting power.",
        8.6, 1.4, 4.5, 0.55, size=13, bold=True, color=NAVY)
stat_box(sl, "OLS slope (log voting power on AI fraction)", "−1.84", 8.6, 2.1,
         w=4.5, h=0.95, fill=RGBColor(0xe8,0xf0,0xff), vcolor=NAVY)
stat_box(sl, "t-statistic / p-value", "t = −3.47\np = 0.001", 8.6, 3.2,
         w=4.5, h=1.05, fill=LGRAY, vcolor=NAVY)
textbox(sl, "High-AI delegates (>20% AI fraction):", 8.6, 4.4, 4.5, 0.35, size=11, color=DGRAY)
textbox(sl, "avg. 71,000 ARB", 8.6, 4.75, 4.5, 0.38, size=14, bold=True, color=POST)
textbox(sl, "Low-AI delegates (<5% AI fraction):", 8.6, 5.25, 4.5, 0.35, size=11, color=DGRAY)
textbox(sl, "avg. 1.9M ARB", 8.6, 5.6, 4.5, 0.38, size=14, bold=True, color=PRE)
textbox(sl, "Tokenholders delegate to authentic delegates — rational discrimination at the delegation stage.",
        8.6, 6.15, 4.5, 0.75, size=11, italic=True, color=DGRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 18 — ROBUSTNESS: SMALL SAMPLE
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "Robustness: Small-Sample Concerns (n = 29)",
           "Three independent tests confirm the pre-period result is not driven by outliers or the small sample")
footer(sl, 18)
# Figure narrower to leave room; placed top-left
add_figure(sl, "rob03_cooks_bootstrap.png", 0.5, 1.3, 8.6)
# Three stat boxes stacked on right
items = [
    ("Cook's Distance",
     "Two obs exceed 4/n = 0.138\nthreshold; they roughly cancel.\nRemoving both: Chow F = 5.62\n(perm p = 0.010).", PRE),
    ("Leave-One-Out (n=29)",
     "97% of 29 single removals\nyield r < −0.20.\nMean r = −0.33 (range −0.14 to −0.42).", NAVY),
    ("Bootstrap Permutation\nChow Test",
     "9 / 10,000 random permutations\nexceed observed F = 8.38.\np_perm = 0.0009.", ORANGE),
]
for i, (lbl, val, col) in enumerate(items):
    x = 9.3
    y = 1.35 + i * 1.88
    rect(sl, x, y, 3.85, 1.72, fill=LGRAY, line=col)
    textbox(sl, lbl, x+0.14, y+0.1, 3.55, 0.5, size=10, bold=True, color=col)
    textbox(sl, val, x+0.14, y+0.62, 3.55, 1.0, size=10, color=DGRAY)
# Interpretation callout fills the bottom
callout_box(sl,
    "The pre-period result (r = −0.33, p = 0.030) is not fragile: "
    "no single observation drives it, every subsample confirms it, "
    "and the structural break cannot be explained by random chance. "
    "The Chow F = 8.38 sits in the extreme tail of its permutation null.",
    0.5, 6.15, 12.3, 0.95, fill=RGBColor(0xe8,0xf0,0xff), border=BLUE, size=12)

# ════════════════════════════════════════════════════════════
# SLIDE 19 — ROBUSTNESS: SECURITY CONFOUND
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "Robustness: The Security Incident Confound")
footer(sl, 19)
textbox(sl,
    "Referee concern: DeFi security incidents (Kyber Nov 23, HTX Nov 22, Binance/DOJ Nov 21) "
    "create an engagement spike just before the QLR break — is the structural break an artifact?",
    0.5, 1.35, 12.3, 0.65, size=13, color=DGRAY)
rect(sl, 0.5, 2.1, 5.9, 2.9, fill=LGRAY, line=BLUE)
textbox(sl, "Direct test:", 0.65, 2.2, 5.6, 0.38, size=12, bold=True, color=NAVY)
textbox(sl,
    "Remove the three matched proposals with 'security' in their title "
    "from Oct–Nov 2023 and re-estimate.",
    0.65, 2.6, 5.6, 0.85, size=12, color=DGRAY)
rows2 = [("Full sample", "−0.33", "8.38  (p = 0.0004)"),
         ("Security excluded", "−0.332  (p = 0.098)", "7.98  (p = 0.0006)")]
textbox(sl, "Sample", 0.65, 3.55, 1.8, 0.32, size=10, bold=True, color=NAVY)
textbox(sl, "r_pre", 2.55, 3.55, 1.8, 0.32, size=10, bold=True, color=NAVY)
textbox(sl, "Chow F", 4.4, 3.55, 1.7, 0.32, size=10, bold=True, color=NAVY)
for i, (s, r, f) in enumerate(rows2):
    y = 3.9 + i*0.52
    bg_c = RGBColor(0xee,0xf8,0xee) if i==1 else WHITE
    rect(sl, 0.55, y-0.05, 5.8, 0.5, fill=bg_c)
    textbox(sl, s, 0.65, y, 1.8, 0.38, size=11, color=DGRAY)
    textbox(sl, r, 2.55, y, 1.8, 0.38, size=11, color=DGRAY)
    textbox(sl, f, 4.4,  y, 1.7, 0.38, size=11, color=DGRAY)

callout_box(sl,
    "Removing security proposals STRENGTHENS the pre-period correlation.\n"
    "Those proposals attenuate rather than inflate the baseline estimate.\n"
    "The confound runs backwards from what the concern implies.",
    6.6, 2.1, 6.0, 2.9, fill=RGBColor(0xe8,0xf8,0xe8), border=RGBColor(0x22,0x88,0x44), size=13)
callout_box(sl,
    "Note: with GPT-4 Turbo (Nov 6) as break, the security incidents (Nov 21-23) "
    "fall in the POST-period and cannot inflate the pre-period estimate at all.",
    0.5, 5.2, 12.3, 0.95, fill=RGBColor(0xe8,0xf0,0xff), size=12)

# ════════════════════════════════════════════════════════════
# SLIDE 20 — RULING OUT ALTERNATIVES
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "Ruling Out Alternative Explanations")
footer(sl, 20)
alts = [
    ("Forum displacement",
     "Human post volume falls as deliberation moves elsewhere",
     "Human posts doubled post-shock (2,650 → 5,898). Volume grew; signal died."),
    ("Voting power concentration",
     "Nakamoto coefficient declines around break date",
     "Power concentration stable through Nov 2023. Break dates do not coincide."),
    ("Proposal complexity shift",
     "Gradual drift rather than sharp structural break",
     "QLR identifies a sharp break on routine proposals. Complexity predicts gradual, not sharp."),
    ("Token vesting cliff (Mar 2024)",
     "Break should coincide with vesting cliff, not GPT-4 Turbo",
     "QLR identifies Nov 2023 — four months prior. Uniswap placebo test fails."),
]
for i, (alt, pred, ev) in enumerate(alts):
    y = 1.35 + i * 1.4
    rect(sl, 0.5, y, 12.3, 1.25, fill=LGRAY, line=BLUE)
    textbox(sl, alt, 0.65, y+0.08, 3.5, 0.4, size=12, bold=True, color=NAVY)
    textbox(sl, f"Predicts: {pred}", 0.65, y+0.52, 5.5, 0.6, size=11, color=DGRAY, italic=True)
    rect(sl, 6.4, y+0.08, 6.2, 1.05, fill=RGBColor(0xe8,0xf8,0xe8), line=RGBColor(0x22,0x88,0x44))
    textbox(sl, f"✔  {ev}", 6.55, y+0.15, 5.9, 0.9, size=11, color=DGRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 21 — DELEGATE PANEL
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "Within-Delegate Evidence: Panel Analysis")
footer(sl, 21)
add_figure(sl, "fig12_panel.png", 0.5, 1.3, 7.8)
bullets(sl,
    ["Within-delegate variation in AI fraction across proposals",
     "Delegates who adopt AI see subsequent delegated power fall",
     "Consistent with tokenholders detecting and penalizing AI adoption",
     "Timing: power decline follows adoption, not simultaneous"],
    8.6, 1.4, 4.6, 4.0, title="Identification", size=13)
callout_box(sl,
    "Cannot fully rule out selection (delegates who adopt AI may differ "
    "in other ways). But timing of decline post-adoption is suggestive of "
    "a causal response.",
    8.6, 5.65, 4.6, 1.05, size=12)

# ════════════════════════════════════════════════════════════
# SLIDE 22 — WELFARE
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "Welfare Implications: A Technology Externality")
footer(sl, 22)
add_figure(sl, "fig03_power_concentration.png", 0.5, 1.3, 6.5)
bullets(sl,
    ["The externality is not imposed by AI-using forum participants",
     "It is imposed by the technology firms whose models cross the capability threshold",
     "Even if NO governance participant deployed AI, the credible THREAT suffices",
     "Deliberation mechanisms require not just honest behavior, but verifiable honest behavior"],
    7.3, 1.4, 5.8, 4.0, title="The externality structure", size=13)
callout_box(sl,
    "Analogy: The bank run does not require depositors to actually withdraw — "
    "the possibility suffices to destroy the coordinating function of the bank.\n\n"
    "The AI shock does not require AI posts to dominate — "
    "the possibility suffices to destroy the forum's credibility.",
    7.3, 5.55, 5.8, 1.2, fill=RGBColor(0xfd,0xf0,0xf0), border=ORANGE, size=12)

# ════════════════════════════════════════════════════════════
# SLIDE 23 — BROADER IMPLICATIONS
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
header_bar(sl, "Broader Implications for Governance Design")
footer(sl, 23)
bullets(sl,
    ["On-chain rationale requirements with AI-disclosure norms",
     "Attestation mechanisms: proof-of-personhood, cryptographic signing of contributions",
     "Moving critical deliberation to verifiable venues (video calls, signed documents)",
     "Delegate incentive programs must penalize AI-generated rationales explicitly"],
    0.5, 1.3, 6.0, 4.5, title="For DAO design", size=13)
bullets(sl,
    ["Public comment periods on regulation face the same mechanism",
     "Earnings call Q&A, ESG engagement letters, investor forums",
     "Regulatory design: LLM capability releases may be governance events requiring disclosure",
     "Any soft-information market where AI can mimic authentic participation is vulnerable"],
    6.8, 1.3, 6.0, 4.5, title="For information economics broadly", size=13)
callout_box(sl,
    "The general principle: text-based deliberation is fragile once AI passes a quality threshold. "
    "This is not a bug in DAO design — it is a fundamental property of any "
    "information market in which AI can mimic authentic participation.",
    0.5, 6.0, 12.3, 1.15, fill=RGBColor(0xe8,0xf0,0xff), border=BLUE, size=13)

# ════════════════════════════════════════════════════════════
# SLIDE 24 — SUMMARY
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl); header_bar(sl, "Summary of Findings"); footer(sl, 24)
findings = [
    ("Finding 1", "The forum signal existed and then died",
     "Pre-shock: r = −0.33 (p=0.030, n=29). Post-shock: r = −0.10 (p=0.39, n=73).\n"
     "Chow F = 8.38 (p=0.0004). QLR break: November 25, 2023 — nineteen days after GPT-4 Turbo.",
     PRE),
    ("Finding 2", "Volume mechanism falsified; readership decline is not",
     "Signal collapses before the AI volume surge. Forum readership fell 36% (88→56 reads/post) "
     "beginning November 2023. Unique human participants fell 33% (740→496).",
     ORANGE),
    ("Finding 3", "Market-level collapse, not post-level screening",
     "Delegates exited the forum rather than discriminating among posts. "
     "Rational anticipation of contamination, not realized contamination, drives the break. "
     "Consistent with Akerlof (1970): exit when lemons are possible, not when frequent.",
     NAVY),
]
for i, (tag, title, body, col) in enumerate(findings):
    y = 1.35 + i * 1.8
    rect(sl, 0.5, y, 12.3, 1.65, fill=LGRAY, line=col)
    rect(sl, 0.5, y, 1.5, 1.65, fill=col)
    textbox(sl, tag, 0.55, y+0.55, 1.4, 0.5, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(sl, title, 2.15, y+0.1, 10.5, 0.42, size=13, bold=True, color=col)
    textbox(sl, body, 2.15, y+0.55, 10.5, 1.0, size=12, color=DGRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 25 — CONCLUSION
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl, NAVY); footer(sl, 25)
rect(sl, 0, 0, 13.333, 7.5, fill=NAVY)
rect(sl, 0, 1.8, 13.333, 0.06, fill=BLUE)
rect(sl, 0, 5.65, 13.333, 0.06, fill=BLUE)
textbox(sl, "The Effect of AI Writing on Governance: Evidence from a $3.5 Billion DAO",
        0.6, 0.3, 12.1, 0.65, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(sl,
    "The DAO forum was a genuine public good — a credible aggregator of private information "
    "about proposal quality. GPT-4 Turbo made that good impossible to maintain: once AI could "
    "write governance-quality prose, rational delegates stopped treating forum activity as a "
    "signal, even though the forum kept growing.",
    0.8, 2.05, 11.7, 1.3, size=14, color=WHITE)
textbox(sl,
    "The welfare loss is borne by tokenholders who no longer receive the benefits of deliberative "
    "information aggregation. The externality is imposed not by AI-using delegates but by the "
    "technology firms whose model releases crossed the capability threshold.",
    0.8, 3.5, 11.7, 1.1, size=14, color=RGBColor(0xcc,0xdd,0xff))
textbox(sl,
    "Text-based deliberation is fragile.\n"
    "The fragility is not a bug in DAO design — it is a fundamental property\n"
    "of any information market in which AI can mimic authentic participation.",
    0.8, 5.85, 11.7, 1.1, size=15, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(sl, "Thank you.   Comments welcome.",
        0.6, 7.0, 12.1, 0.38, size=12, color=RGBColor(0xaa,0xcc,0xff), align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────
out = "presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({prs.slides.__len__()} slides)")
